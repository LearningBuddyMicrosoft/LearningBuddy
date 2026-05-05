import os
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image

class DocumentChunker:
    def __init__(self, overlap_size=25, chunk_size=50):
        self.overlap_size = overlap_size
        self.chunk_size = chunk_size

    def chunk_pdf(self, file_path):
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(reader.pages, start=1):
                    extracted = page.extract_text()
                    if extracted:
                        page_text = f"Page {page_num}: {extracted}"
                        # Chunk the page text
                        page_chunks = self._chunk_text(page_text)
                        chunks.extend(page_chunks)
                
                # 🌟 THE SANITY CHECK 🌟
                print(f"📄 Extracted a total of {len(chunks)} chunks from the PDF.")
        except Exception as e:
            print(f"Error processing PDF: {e}")
        return chunks

    def chunk_docx(self, file_path):
        chunks = []
        try:
            doc = Document(file_path)
            text = ' '.join([para.text for para in doc.paragraphs])
            chunks.extend(self._chunk_text(text))
        except Exception as e:
            print(f"Error processing DOCX: {e}")
        return chunks

    def chunk_pptx(self, file_path):
        chunks = []
        try:
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
            chunks.extend(self._chunk_text(text))
        except Exception as e:
            print(f"Error processing PPTX: {e}")
        return chunks

    def chunk_image(self, file_path):
        chunks = []
        try:
            # Assuming OCR extraction from image; placeholder for actual implementation
            text = "Extracted text from image using OCR"
            chunks.extend(self._chunk_text(text))
        except Exception as e:
            print(f"Error processing image: {e}")
        return chunks

    def _chunk_text(self, text):
        # Token-aware semantic chunking
        tokens = text.split()
        chunks = []
        for i in range(0, len(tokens), self.chunk_size - self.overlap_size):
            chunk = ' '.join(tokens[i:i + self.chunk_size])
            if chunk: 
                chunks.append(chunk)
        return chunks

    def chunk_file(self, file_path):
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf':
            return self.chunk_pdf(file_path)
        elif extension == '.docx':
            return self.chunk_docx(file_path)
        elif extension == '.pptx':
            return self.chunk_pptx(file_path)
        elif extension in ['.jpg', '.jpeg', '.png']:
            return self.chunk_image(file_path)
        else:
            raise ValueError("Unsupported file type")

# Example usage
if __name__ == '__main__':
    chunker = DocumentChunker()
    file_path = 'example.pdf'  # Replace with actual file path
    chunks = chunker.chunk_file(file_path)
    print(chunks)
